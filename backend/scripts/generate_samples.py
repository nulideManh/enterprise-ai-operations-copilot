import os
from docx import Document
from pptx import Presentation

def create_samples():
    output_dir = os.path.abspath(os.path.join(os.path.dirname(__file__), "../../sample_documents"))
    os.makedirs(output_dir, exist_ok=True)
    
    # 1. Create hr_policy.docx
    doc1 = Document()
    doc1.add_heading("Sổ tay nhân sự và chính sách HR", 0)
    doc1.add_heading("1. Giờ làm việc cốt lõi", 1)
    doc1.add_paragraph(
        "Giờ làm việc cốt lõi của công ty là từ 9:00 đến 17:00, từ thứ Hai đến thứ Sáu. "
        "Nhân viên cần sẵn sàng phối hợp trong khung giờ này. Các thỏa thuận làm việc linh hoạt "
        "có thể được trao đổi với trưởng bộ phận."
    )
    doc1.add_heading("2. Nghỉ phép năm và thời gian nghỉ", 1)
    doc1.add_paragraph(
        "Tất cả nhân viên toàn thời gian được hưởng 15 ngày nghỉ phép có lương trong mỗi năm dương lịch. "
        "Yêu cầu nghỉ phép phải được gửi qua cổng HR và được quản lý bộ phận phê duyệt "
        "ít nhất 5 ngày làm việc trước ngày bắt đầu nghỉ. Tối đa 5 ngày phép chưa sử dụng "
        "có thể được chuyển sang năm tiếp theo."
    )
    doc1.add_heading("3. Phúc lợi sức khỏe và chăm sóc nhân viên", 1)
    doc1.add_paragraph(
        "Công ty cung cấp bảo hiểm sức khỏe, nha khoa và thị lực cho toàn bộ nhân viên toàn thời gian. "
        "Chi phí bảo hiểm của nhân viên được công ty chi trả 100%, và chi trả 50% cho người phụ thuộc. "
        "Ngoài ra, nhân viên nhận khoản hỗ trợ chăm sóc sức khỏe 50 USD mỗi tháng cho phòng tập hoặc hoạt động wellness."
    )
    doc1.save(os.path.join(output_dir, "hr_policy.docx"))
    print("Created hr_policy.docx")

    # 2. Create it_security_guide.pptx
    prs = Presentation()
    
    # Slide 1: Title
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Chính sách bảo mật CNTT và làm việc từ xa"
    subtitle.text = "Hướng dẫn vận hành an toàn"

    # Slide 2: VPN
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "Quy định VPN và kết nối từ xa"
    tf = body_shape.text_frame
    tf.text = "Luôn sử dụng ứng dụng Cisco VPN chính thức của công ty khi làm việc từ xa."
    p2 = tf.add_paragraph()
    p2.text = "Không truy cập tài nguyên công ty từ Wi-Fi công cộng nếu chưa bật VPN."
    p3 = tf.add_paragraph()
    p3.text = "Xác thực đa yếu tố (MFA) qua Duo Security là bắt buộc cho mọi kết nối VPN."

    # Slide 3: Passwords
    slide = prs.slides.add_slide(slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "Yêu cầu độ mạnh mật khẩu"
    tf = body_shape.text_frame
    tf.text = "Mật khẩu phải có tối thiểu 12 ký tự."
    p2 = tf.add_paragraph()
    p2.text = "Mật khẩu phải có chữ hoa, chữ thường, chữ số và ký tự đặc biệt như !, @, #."
    p3 = tf.add_paragraph()
    p3.text = "Mật khẩu phải được thay đổi mỗi 90 ngày và không được dùng lại 5 mật khẩu gần nhất."
    
    prs.save(os.path.join(output_dir, "it_security_guide.pptx"))
    print("Created it_security_guide.pptx")

    # 3. Create finance_instructions.docx
    doc2 = Document()
    doc2.add_heading("Chính sách chi tiêu và phê duyệt chi phí", 0)
    doc2.add_heading("1. Nguyên tắc chi tiêu chung", 1)
    doc2.add_paragraph(
        "Mọi khoản chi của công ty phải phục vụ mục đích công việc, có hóa đơn hoặc chứng từ chi tiết, "
        "và được nộp trong vòng 30 ngày kể từ ngày giao dịch. Giao dịch trên 25 USD không có chứng từ sẽ không được hoàn tiền."
    )
    doc2.add_heading("2. Ma trận hạn mức phê duyệt", 1)
    p = doc2.add_paragraph()
    p.add_run("Các hạn mức sau áp dụng cho mọi khoản mua sắm của bộ phận:\n")
    p.add_run("- Chi phí dưới 500 USD: cần trưởng nhóm phê duyệt.\n")
    p.add_run("- Chi phí từ 500 đến 5.000 USD: cần quản lý bộ phận phê duyệt.\n")
    p.add_run("- Chi phí trên 5.000 USD: cần CFO hoặc giám đốc điều hành phê duyệt.\n")
    doc2.add_heading("3. Quy trình gửi hóa đơn", 1)
    doc2.add_paragraph(
        "Hóa đơn từ nhà cung cấp phải được gửi trực tiếp đến finance-receipts@enterprise.com. "
        "Mỗi hóa đơn cần có tên nhà cung cấp, số hóa đơn, số tiền, loại tiền tệ và ngày phát hành. "
        "Thanh toán được xử lý theo lịch net-30 sau khi bộ phận liên quan xác nhận lần cuối."
    )
    doc2.save(os.path.join(output_dir, "finance_instructions.docx"))
    print("Created finance_instructions.docx")

if __name__ == "__main__":
    create_samples()
