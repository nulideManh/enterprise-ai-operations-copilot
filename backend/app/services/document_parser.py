from io import BytesIO
from pathlib import Path
import re

from docx import Document as DocxDocument
from pptx import Presentation
from pypdf import PdfReader


def parse_document(filename: str, content: bytes) -> list[tuple[int | None, str]]:
    suffix = Path(filename).suffix.lower()
    if suffix == ".pdf":
        reader = PdfReader(BytesIO(content))
        pages: list[tuple[int | None, str]] = []
        for index, page in enumerate(reader.pages):
            text = page.extract_text() or ""
            if text.strip():
                pages.append((index + 1, text))
        return pages

    if suffix == ".docx":
        doc = DocxDocument(BytesIO(content))
        text = "\n".join(paragraph.text for paragraph in doc.paragraphs if paragraph.text.strip())
        return [(None, text)] if text.strip() else []

    if suffix == ".pptx":
        deck = Presentation(BytesIO(content))
        pages: list[tuple[int | None, str]] = []
        for index, slide in enumerate(deck.slides):
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text)
            if texts:
                pages.append((index + 1, "\n".join(texts)))
        return pages

    raise ValueError("Định dạng tài liệu chưa được hỗ trợ. Vui lòng tải lên PDF, DOCX hoặc PPTX.")


def recursive_chunk(text: str, max_chars: int = 1200, overlap: int = 180) -> list[str]:
    clean = " ".join(text.split())
    if not clean:
        return []
    chunks: list[str] = []
    start = 0
    while start < len(clean):
        end = min(start + max_chars, len(clean))
        if end < len(clean):
            split_at = max(clean.rfind(". ", start, end), clean.rfind("\n", start, end))
            if split_at > start + max_chars // 2:
                end = split_at + 1
        chunks.append(clean[start:end].strip())
        if end == len(clean):
            break
        start = max(0, end - overlap)
    return chunks


def semantic_chunk(text: str, max_chars: int = 1400, min_chars: int = 260, overlap_sentences: int = 1) -> list[str]:
    clean = "\n".join(line.strip() for line in text.splitlines() if line.strip())
    if not clean:
        return []

    paragraphs = [part.strip() for part in re.split(r"\n{2,}", clean) if part.strip()]
    units: list[str] = []
    for paragraph in paragraphs or [clean]:
        if len(paragraph) <= max_chars:
            units.append(paragraph)
            continue
        sentences = [sentence.strip() for sentence in re.split(r"(?<=[.!?])\s+", paragraph) if sentence.strip()]
        units.extend(sentences or [paragraph])

    chunks: list[str] = []
    current: list[str] = []
    current_len = 0

    for unit in units:
        unit_len = len(unit)
        if current and current_len + unit_len + 2 > max_chars and current_len >= min_chars:
            chunks.append("\n\n".join(current).strip())
            current = current[-overlap_sentences:] if overlap_sentences > 0 else []
            current_len = sum(len(item) + 2 for item in current)
        current.append(unit)
        current_len += unit_len + 2

    if current:
        chunks.append("\n\n".join(current).strip())

    return [chunk for chunk in chunks if chunk]


def chunk_text(text: str, strategy: str = "recursive") -> list[str]:
    if strategy == "semantic":
        return semantic_chunk(text)
    return recursive_chunk(text)
